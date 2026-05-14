# app/pipeline/ingest.py

import os
from datetime import datetime
from typing import List
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    UnstructuredFileLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from app.config import settings
from app.core.logger import logger

# 增加office 文檔處理能力，放在函數内加載
# import pandas as pd
# from docx import Document as DocxDocument
# from pptx import Presentation

# --- 切片器实例化外移，减少重复开销 ---
_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=settings.CHUNK_SIZE,
    chunk_overlap=settings.CHUNK_OVERLAP
)

def list_all_files(directory: str) -> List[str]:
    """
    【原子接口 1】：递归文件扫描器。
    职责：仅负责物理扫描，为多线程提供任务清单。
    """
    file_list = []

    if not os.path.exists(directory):
        logger.error(f"❌ 根目录不存在: {directory}")
        return []

    for root, dirs, files in os.walk(directory):
        for filename in files:
            # 过滤临时文件和隐藏文件
            if filename.startswith(('.', '~')):
                continue
            file_list.append(os.path.join(root, filename))

    return file_list


def process_file_to_docs(file_path: str) -> List[Document]:
    """
    【原子接口 2】：单文件并发处理器（Task B 核心）。
    """
    filename = os.path.basename(file_path)
    base_directory = settings.DATA_UPLOAD_DIR
    ingest_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    rel_path = os.path.relpath(os.path.dirname(file_path), base_directory)
    domain = "未分类资产" if rel_path == "." else rel_path

    try:
        # 1. 动态选择加载器
        if filename.lower().endswith(".pdf"):
            loader = PyPDFLoader(file_path)
            raw_docs = loader.load()  # 保持原樣執行
        elif filename.lower().endswith(".txt"):
            # --- [修正：建立中文优先的编码探测逻辑] ---
            # 放弃直接调用 autodetect_encoding，改为手动尝试中文编码链
            raw_docs = []
            encodings_to_try = ["GB2312","utf-8", "GBK", "GB18030","big5"]
            success = False

            for enc in encodings_to_try:
                try:
                    # 原地尝试加载，失败则捕获异常进入下一循环
                    loader = TextLoader(file_path, encoding=enc)
                    raw_docs = loader.load()
                    success = True
                    logger.info(f"📄 [編碼成功] {filename} 已使用 {enc} 讀取")
                    break
                except Exception:
                    continue

            if not success:
                # 如果中文链全灭，最后报错拦截，绝不回退到 latin-1
                raise ValueError(f"無法識別的文本編碼，請將文件轉為 UTF-8")

        # ========== 新增：Word 文件支援 ==========
        elif filename.lower().endswith(".docx"):
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            content = "\n".join(full_text)
            # 手動建立 Document 物件
            raw_docs = [Document(page_content=content, metadata={"source": file_path})]
            logger.info(f"📄 [Word] {filename} 讀取成功，共 {len(full_text)} 段落")

        # ========== 新增：Excel 文件支援 ,當前暫緩處理==========
        #elif filename.lower().endswith((".xlsx", ".xls")):
        #    import pandas as pd
        #    df = pd.read_excel(file_path, engine='openpyxl')
        #    # 將 DataFrame 轉為文字
        #    content = df.to_string()
        #    from langchain_core.documents import Document
        #    raw_docs = [Document(page_content=content, metadata={"source": file_path})]
        #    logger.info(f"📄 [Excel] {filename} 讀取成功，共 {len(df)} 行")
        # ========== 新增：Excel 文件支援 ,當前暫緩處理==========

        # =============拒絕excel+++++++++++++++ #
        elif filename.lower().endswith((".xlsx", ".xls")):
            raise NotImplementedError(
                f"Excel 文件 {filename} 當前處於實驗性支持階段，"
                "存在解析穩定性問題。暫時請轉換為 CSV 或 TXT 格式後導入。"
            )

        # ========== 新增：PowerPoint 文件支援 ==========
        elif filename.lower().endswith((".pptx", ".ppt")):
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            content = "\n".join(full_text)
            raw_docs = [Document(page_content=content, metadata={"source": file_path})]
            logger.info(f"📄 [PowerPoint] {filename} 讀取成功，共 {len(full_text)} 個文本元素")
        else:
            loader = UnstructuredFileLoader(file_path)
            raw_docs = loader.load()

        # 2. 注入 Metadata 契约 (歸一化路徑標籤)

        # 獲取相對於數據根目錄的相對路徑，並統一轉換為正斜槓，解決跨平台對齊問題
        # 將 'folder\file.txt' 統一轉為 'folder/file.txt'，作為全局唯一識別碼
        rel_file_path = os.path.relpath(file_path, base_directory)
        normalized_source = rel_file_path.replace('\\', '/')

        for doc in raw_docs:
            doc.metadata["source"] = normalized_source
            doc.metadata["domain"] = domain
            doc.metadata["ingest_time"] = ingest_time

        # 3. 执行物理切片
        splits = _SPLITTER.split_documents(raw_docs)
        logger.info(f"✅ [已标记 - {domain}] 加载成功: {filename} ({len(splits)} chunks)")
        return splits

    except Exception as e:
        logger.error(f"❌ 加载失败 [{filename}]: {str(e)}")
        return []


def load_and_split_all_documents() -> List[Document]:
    """
    【向后兼容接口】：
    内部通过原子函数串行实现，但在 ingest_service 中会被 _parallel_load_and_split 替代。
    """
    files = list_all_files(settings.DATA_UPLOAD_DIR)
    all_documents = []

    for file_path in files:
        all_documents.extend(process_file_to_docs(file_path))

    return all_documents

def ingest_documents():
    """
    【向后兼容接口】：
    为了不破坏原有的逻辑调用，保留此入口，但内部直接调用 Service 层。
    体现了“旧契约不变，新引擎升级”的工程平替思路。
    """
    from app.services.ingest_service import initialize_knowledge_base
    logger.info("🔄 [系统路由] 正在通过向后兼容接口触发全量重塑...")
    return initialize_knowledge_base(force_rebuild=True)